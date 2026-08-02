#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate per-project MDs, integrated MD, and portfolio PPTX for 이지우."""

from __future__ import annotations

import textwrap
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
PROJECTS_DIR = ROOT / "projects"
PROJECTS_DIR.mkdir(parents=True, exist_ok=True)

# Design tokens (KB ink + iM teal + WJVOX calm neutrals)
INK = RGBColor(0x26, 0x28, 0x2C)
MUTED = RGBColor(0x5C, 0x63, 0x6B)
ACCENT = RGBColor(0x0D, 0x94, 0x88)  # teal
ACCENT_SOFT = RGBColor(0xCC, 0xF2, 0xEC)
BG = RGBColor(0xF5, 0xF6, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD8, 0xDE, 0xE4)
WARN = RGBColor(0xB0, 0x7D, 0x00)

FONT = "Pretendard"
FONT_FALLBACK = "AppleGothic"

# --- Content inventory -------------------------------------------------------

PROFILE = {
    "name": "이지우",
    "title": "AI가 답을 던지는 시대에 왜?를 묻는 개발자 이지우입니다",
    "school": "광운대학교 인공지능융합대학 정보융합학부 비주얼테크놀로지전공",
    "github": "https://github.com/dlwldn4824",
    "instagram": "@due_study_archive",
    "admission": "2024.03",
}

GRADES = {
    "gpa": "4.32 / 4.5",
    "percentile": "97.9",
    "credits": "101",
    "semesters": [
        ("2024-1", "3.83", "프로그래밍입문·통계학입문 A+, 미적분학1 A0"),
        ("2024-2", "4.50", "AI와컴퓨팅사고·이산수학·대학영어 A+ (만점 학기)"),
        ("2025-1", "4.41", "빅데이터프로그래밍·AI수학·UX/UI·OOP A+"),
        ("2025-2", "4.41", "자료구조·정보디자인프로그래밍실습·HCI와UX평가 A+"),
        ("2026-1", "4.50", "오픈소스실습·기계학습·DB·텍스트마이닝·모바일로봇 A+ (만점 학기)"),
    ],
    "scholarships": [
        "2024.09 성적우수 장학금 반액 2,236,000원",
        "2025.06 / 2025.09 성적우수 장학금 (쿼터)",
        "2025.12 성적우수 장학금 쿼터 1,148,250원",
        "2026 호반장학생(이노베이션) 지원 — 인프라·인터뷰·교육 활용 계획",
    ],
    "awards": [
        "2024.05 글쓰기 대회 가작(산문) — 광운대",
        "2024.09·12 창업 동아리 경진대회 장려상 — 광운대",
        "2024.12 Dean’s List — 광운대",
        "2025.05.29 매치업 심화과정 경진대회 우수상 — 광운대",
        "2025.12 마이크로모듈 초급 SS급 — 지능형로봇사업단",
        "2026.06.26 지능형 로봇 컨소시엄 대상 → Smart Icheon Care",
        "2026 보조공학 장려상 → HOPE",
        "2026.07.03 HUSS AI 경진대회 장려상 → HOPE (로컬 임팩트 표기 삭제)",
    ],
}

RESEARCH = [
    {
        "slug": "research-hope-validation",
        "name": "HOPE 문제 검증 연구",
        "role": "연구 설계·텍스트마이닝 파이프라인",
        "why": "조음음운장애 아동 언어재활에서 ‘가정 피드백 공백’이 실제 Pain인지 데이터로 검증하기 위해 선택",
        "validation": "공공데이터·뉴스·논문 → TF-IDF / LDA / 공출현 네트워크로 Pain Point 도출 후 기능 매핑",
        "solution": "재현 가능한 Python 파이프라인(수집→전처리→토픽→교차비교)과 synthesis 리포트",
        "impact": "보조공학 장려상 + HUSS AI 경진대회 장려상의 문제 정의 근거. HUSS 로컬 임팩트 표기 미사용",
        "strengths": "문제 정의, 연구 설계, 텍스트마이닝, 근거 기반 제품 기획",
        "stack": "Python, KoNLPy, TF-IDF, LDA, NetworkX, matplotlib",
        "github": "https://github.com/dlwldn4824/HOPE_organization",
    },
    {
        "slug": "research-mental-health",
        "name": "Multi-Agent Mental Health Evaluation",
        "role": "연구용 평가 프레임 · RAG/Eval 파트",
        "why": "단일 LLM 상담 응답의 증상·위험·안전 추론 품질을 객관적으로 비교할 필요가 있어 선택",
        "validation": "MentalChat16K 기반, BERTScore/ROUGE·Judge·ChromaDB 공식 KB(NIMH/WHO/NICE)",
        "solution": "Symptom→Risk→Safety→Consensus 다층 에이전트 + RAG 대비 Single LLM 벤치마크",
        "impact": "P/R/F1, high-risk recall, unsafe rate 등 실험 지표로 ‘안전한 AI’ 설계 근거 확보",
        "strengths": "연구 실험 설계, RAG, 안전성 평가, 멀티에이전트 오케스트레이션",
        "stack": "Python, Ollama(qwen2.5:7b), ChromaDB, BERTScore/ROUGE",
        "github": "https://github.com/dlwldn4824/TM-MultiLayer-MentalHealth",
    },
    {
        "slug": "research-ml-dessert",
        "name": "디저트 소비 선호도 ML 예측",
        "role": "데이터 분석·모델 비교 (기계학습)",
        "why": "인플레이션 환경에서 소비 변화는 ‘감’이 아니라 행정동 단위 데이터로 검증해야 한다고 판단",
        "validation": "서울시 상권분석 서비스 5개년 · 카페·제과점 필터 · 행정동 단위",
        "solution": "전처리→탐색→복수 ML 모델 학습·비교 파이프라인 (notebooks + src)",
        "impact": "실데이터 기반 예측 모델 비교 경험. GPA 만점 학기(기계학습 A+)와 연결되는 실습 산출물",
        "strengths": "데이터 전처리, 모델 비교, 도메인 해석",
        "stack": "Python, scikit-learn, pandas, Jupyter",
        "github": "https://github.com/dlwldn4824/machine_learning",
    },
]

PROJECTS = [
    {
        "slug": "hope",
        "name": "HOPE",
        "tagline": "보조공학·HUSS AI 이중 장려상 · 아동 조음 AI 보조",
        "tier": "S",
        "why": "치료실 ~2시간(주 1%) 대비 가정 166시간은 객관적 피드백이 거의 없는 구조적 공백. 보조공학(가정 재활 보조) + HUSS(사회문제×AI)로 동시에 설득되는 문제로 선택",
        "validation": "뉴스·논문·공공데이터 TF-IDF/LDA/네트워크. Pain(비용·대기·가정피드백·일반화) 도출. HUSS 심사 후 최우수팀 인터뷰·기준 재분석. HUSS 로컬 임팩트 표기는 사용하지 않음",
        "solution": "게임형 연습 + Speech Coach/PCC + 보호자 리포트. 치료사 대체가 아닌 가정 보조 도구",
        "impact": "① 보조공학 장려상 ② HUSS AI 경진대회 장려상(2026.07.03). 레포: github.com/dlwldn4824/HOPE_organization",
        "role": "팀장(정융고능팀) — UI/UX·게임 인터페이스·콘텐츠 총괄",
        "strengths": "리더십, 문제정의, 보조공학 UX, 연구→제품 번역",
        "stack": "React, TypeScript, Vite, Tailwind v4, Node, Speech Coach API",
        "github": "https://github.com/dlwldn4824/HOPE_organization",
        "extra": "이중 장려상 · 로컬 임팩트 미표기",
    },
    {
        "slug": "im-ready",
        "name": "iM Ready / IM : Shield",
        "tagline": "피싱 전 AI 백신 훈련 + 후 골든타임 대응",
        "tier": "S",
        "why": "2025 보이스피싱 피해 약 1조 2,578억. 차단 기술은 늘어도 피해는 늘고, 피해자 65.3%는 ‘판단할 틈 박탈’ — 예방·사후 공백을 메우기 위해 선택",
        "validation": "금감원·경찰청·경기도 설문(1,195명) 등 사회통계·규제 프레임 검토",
        "solution": "모의 피싱 훈련(행동로그만) + 골든타임(공식 절차 RAG/규칙엔진 우선) + 지연이체·가족공유 안내. 데모: im-ready-fawn.vercel.app",
        "impact": "예방(훈련)과 사후(골든타임)를 한 제품 내러티브로 연결. We are 개발자 피치덱 29p",
        "role": "프론트엔드 핵심 구현 (브랜치 이지우, 피싱백신·골든타임 플로우)",
        "strengths": "모바일 UX, 안전 설계(원본 미수집), 규제 감각, 스토리텔링",
        "stack": "React 19, Vite, TypeScript, Tailwind 4, Pretendard+Jua",
        "github": "https://github.com/ik-s/iM-Ready",
        "extra": "실금융 기능 미연결 데모 — 안전 원칙 명시",
    },
    {
        "slug": "smart-icheon-care",
        "name": "Smart Icheon Care",
        "tagline": "지능형 로봇 컨소시엄 대상 · 지자체 CV 대시보드",
        "tier": "S",
        "why": "불법 현수막 등 현장 순찰 한계. 공개 illegal/legal 2클래스 부재 → HITL로 행정 워크플로에 맞게 선택",
        "validation": "AI Hub·공통테스트 1,892장. 2026 지능형 로봇 컨소시엄 창의융합캠프(후레쉬베리)",
        "solution": "YOLO11s→Risk→클릭 OCR→공무원 CONFIRMED. VWorld GIS·주차·시민신고",
        "impact": "지능형 로봇 컨소시엄 대상(2026.06.26). val mAP50 0.510 / test F1 0.591 mAP50 0.439 / ~15.7 FPS. 레포: smart_icheon_care",
        "role": "주요 기여자 — 풀스택 + CV 파이프라인",
        "strengths": "CV/MLOps, 풀스택, 공공데이터, HITL 설계",
        "stack": "Next.js 16, React 19, TS, Tailwind 4, Leaflet/VWorld, FastAPI, YOLO11s, ByteTrack",
        "github": "https://github.com/dlwldn4824/smart_icheon_care",
        "extra": "대상 수상 ↔ 이 레포 직결",
    },
    {
        "slug": "answer-registry",
        "name": "답변등기 (KB AI Challenge)",
        "tagline": "은행 AI 상담 초안 문장단위 승인·봉인·발송 콘솔",
        "tier": "S",
        "why": "AI 초안 위험 vs 전건 사람승인 부담. ‘승인≠발송’ 사고 위험을 구조적으로 막아야 한다고 판단",
        "validation": "합성 240건 / 실상담 AI Hub 표본 / vitest 불변조건 60건",
        "solution": "위험큐(1,248→깊게 볼 39건) + 정본대조 + 사유정합 + SHA256 digest/HMAC 봉인 + 등기 타임라인. 원칙: 되돌릴 수 있는 곳엔 AI, 없는 곳엔 결정론",
        "impact": "게이트 차단·재생 100%. 큐 3.1%(39/1248). 합성 Recall@39 81.3%→실상담 25.0%(정직하게 격차 공개)",
        "role": "팀 삼삼오오 — UI/디자인 폴리시·구현",
        "strengths": "신뢰성 엔지니어링, 브랜드 UI, 평가 정직성, WebCrypto",
        "stack": "Next.js, Vitest, WebCrypto, KB Yellow #FFCC00, Pretendard, GFC Red Spirit",
        "github": "https://github.com/dlwldn4824/kb_AI_challenge",
        "extra": "데모 https://dlwldn4824.github.io/kb_AI_challenge/",
    },
    {
        "slug": "wjvox",
        "name": "WJVOX (YourVoice)",
        "tagline": "음성 업로드→학습→공개 infer·공유 플랫폼",
        "tier": "A",
        "why": "개인/권한 음성 기반 AI 보이스 생성·공유에서 운영·보안 경계가 실제 병목이라고 판단",
        "validation": "Owner UI 체크포인트, API smoke, RunPod on-demand/GPU 문서화로 운영 검증",
        "solution": "Upload→Train→Publish→Infer. JWT vs WORKER_TOKEN, Cloudflare R2, RunPod GPU worker(RVC 등)",
        "impact": "실서비스 인프라 운영 경험. 3–6개월 운영비 약 100–200만 원 규모 체감(자소서)",
        "role": "UI/UX·타이포·프론트 기여 (메인 커밋 협업)",
        "strengths": "서비스 UX, 보안 경계, GPU 워커 운영 감각",
        "stack": "Next.js, Supabase, Cloudflare R2, RunPod, Vercel, Gumi Romance+Pretendard",
        "github": "https://github.com/KWwoojin/project",
        "extra": "accent #8E60F6, UNIVERSAL_UI 원칙 적용",
    },
    {
        "slug": "bandy",
        "name": "bandy (밴드 예약·공연 관리)",
        "tagline": "게스트 예약·티켓·셋리스트·체크인 — dlwldn4824/band",
        "tier": "A",
        "why": "밴드 공연의 예약(게스트 명단)·티켓·셋리스트·체크인이 엑셀·카톡으로 파편화 — 동아리 실사용 문제로 선택",
        "validation": "노을 밴드 동아리 운영진·공연 현장에서 실사용 플로우로 검증",
        "solution": "관리자 엑셀 업로드 → 게스트 이름·전화 로그인 → 티켓/셋리스트/공연진/이벤트 대시보드",
        "impact": "커밋 320+ 주 기여자. 예약·체크인 단일 웹화. 레포: github.com/dlwldn4824/band",
        "role": "주 기여자 / 풀스택 · 노을 밴드 동아리 운영진",
        "strengths": "실사용 제품화, 도메인 UX, ownership",
        "stack": "React 18, TypeScript, Vite, React Router, xlsx, Firebase, Socket.io, QR",
        "github": "https://github.com/dlwldn4824/band",
        "extra": "인간성·동아리 섹션 핵심 산출물 (조·MT 표기 미사용)",
    },
    {
        "slug": "setlist-recommender",
        "name": "Band Setlist Recommender",
        "tagline": "552곡 기반 커버곡·믹싱 방향 추천",
        "tier": "A",
        "why": "음역·세션·악기·분위기에 맞는 셋리스트·믹스 팁이 없어 리허설 비용이 큼",
        "validation": "오픈소스 실습 기말 — Docker Compose·EC2 배포 체크리스트로 재현성 검증",
        "solution": "rule-based Top5 + 곡별 믹싱 가이드. Streamlit UI + FastAPI + Docker",
        "impact": "552곡 데이터셋. AWS EC2 배포 가능한 교육용 프로덕션 형태",
        "role": "개인 저장소 오너",
        "strengths": "규칙 기반 추천, Docker/EC2, 도메인(밴드) 지식",
        "stack": "Streamlit, FastAPI, Docker Compose, AWS EC2",
        "github": "https://github.com/dlwldn4824/opensource_final",
        "extra": "관련: open-source-fastapi-docker CRUD 실습",
    },
    {
        "slug": "temi-hci",
        "name": "Temi-Tell-Me / TEMI 스탬프 랠리",
        "tagline": "전시장 길안내·줄서기·체험 큐레이션",
        "tier": "A",
        "why": "전시장에서 길 찾기·대기·콘텐츠 선택이 분산되어 방문 경험이 깨짐",
        "validation": "CO-SHOW·HCI 수업/여비 신청 맥락의 현장형 과제",
        "solution": "TEMI 길안내(18존), QR 줄서기, 필터 포토, 퀴즈·추천. Spring Boot + Claude API 이벤트 추천",
        "impact": "18존 내비, 8개 줄서기 프로그램 등 현장 스케일 설계",
        "role": "HCI 저장소 오너 · 테미 팀 협업(yyeonseoo/mobile-robot)",
        "strengths": "HCI/UX 평가, 로봇·모바일 연동, 현장 시나리오",
        "stack": "React, Spring Boot, Capacitor/Android, Claude API, Vite, Flutter(temi)",
        "github": "https://github.com/dlwldn4824/TemiTellMe · https://github.com/dlwldn4824/HCI-UX · https://github.com/dlwldn4824/mobile_robot_temi",
        "extra": "관련: coshow, HCI 저장소",
    },
    {
        "slug": "pintime",
        "name": "PinTime",
        "tagline": "캘린더 기반 일정 조율 데모",
        "tier": "B",
        "why": "주간(시간)·월간(여행) 조율과 공유 링크 UX가 분리되어 있어 하나로 묶고자 함",
        "validation": "서버리스 공유링크 시연 흐름으로 검증",
        "solution": "주간 드래그(시간)/월간(여행), 방·링크, busy 마스킹",
        "impact": "서버 없이 링크에 방 데이터를 실어 기기 간 시연 가능",
        "role": "개인 주도 데모",
        "strengths": "캘린더 UX, 인터랙션 디테일, 경량 아키텍처",
        "stack": "React 19, Vite, Tailwind 4, Lucide, Manrope",
        "github": "(로컬 대외활동/캘린더)",
        "extra": "pin #3b82f6, sidebar #15181f",
    },
    {
        "slug": "cam-kit",
        "name": "Cam-Kit",
        "tagline": "캠퍼스 하이퍼로컬 소분 공동구매 + 스마트 보관함",
        "tier": "B",
        "why": "1인 가구의 대용량·편의점 고가·배달비·수령 시간 불일치",
        "validation": "창업캡스톤 기획 — 타겟·수익모델·하루 5회차 회전 가정",
        "solution": "수요조사→모집결제→소분→보관함 수령",
        "impact": "편의점 대비 최대 60% 할인(기획 지표). 창업 동아리 경진 장려상 맥락과 연결",
        "role": "팀 프로젝트 참여 (remote bhw119/Cam-Kit)",
        "strengths": "비즈니스 모델 설계, 하이퍼로컬 기획, 풀스택 스캐폴드",
        "stack": "Node/Express, MongoDB, JWT, React+Vite, Zustand",
        "github": "https://github.com/bhw119/Cam-Kit",
        "extra": "창업캡스톤",
    },
    {
        "slug": "refund-ranger",
        "name": "환불원정대 (Refund Ranger)",
        "tagline": "티몬·위메프 환불 사태 뉴스 분석 대시보드",
        "tier": "B",
        "why": "환불·민원 뉴스가 파편화되어 여론·토픽을 직관적으로 탐색하기 어려움",
        "validation": "키워드·감정·토픽 분석 파이프라인",
        "solution": "네트워크/차트 + AI 챗봇 「민심이」(Gemini)",
        "impact": "시사 이슈를 데이터 시각화·대화형으로 탐색하는 팀 산출물",
        "role": "팀장",
        "strengths": "리더십, 텍스트분석, 대시보드 UX",
        "stack": "React+Vite, Flask, Gemini, Docker, Nginx, GitHub Actions",
        "github": "(로컬 2학년 2학기/정딥/25-team-refund-ranger)",
        "extra": "정딥 팀 프로젝트",
    },
    {
        "slug": "creative-umc",
        "name": "창의설계 / UMC / 기타 실습",
        "tagline": "기초·커뮤니티·오픈소스 실습 묶음",
        "tier": "B",
        "why": "기초 역량·커뮤니티 학습 루프를 끊지 않기 위해 지속적으로 수행",
        "validation": "수업 산출물·실습 체크리스트·배포",
        "solution": "creative_py_project(Leaflet), UMC 9th_web 포크 작업장, FastAPI Docker CRUD, Open Source Practice/Quiz",
        "impact": "GPA·오픈소스 A+와 연결된 실행 습관. 인스타 회고(@due_study_archive)로 공유",
        "role": "개인·커뮤니티 학습자",
        "strengths": "학습 지속성, 오픈소스 실무 감각",
        "stack": "React, Leaflet, FastAPI, Docker, JavaScript",
        "github": "https://github.com/dlwldn4824/creative_py_project 외",
        "extra": "fork/실습 레포 포함 — 포트폴리오 완결성용",
    },
]

CLUBS = [
    {
        "name": "NOEUL(노을) 밴드 동아리",
        "role": "운영진",
        "points": [
            "공연 예약·체크인·셋리스트를 bandy 웹으로 통합 운영",
            "동아리 Pain을 제품으로 푸는 태도 (내부 조·MT 표기 미사용)",
        ],
    },
    {
        "name": "창업 동아리 / 창업캡스톤",
        "role": "팀원 · 경진 장려상",
        "points": [
            "Cam-Kit 등 비즈니스·제품 동시 설계",
            "장려상 이후에도 회고·개선 루프 유지",
        ],
    },
]

LEADERSHIP = [
    {
        "name": "HOPE / 정융고능팀 팀장",
        "points": [
            "UI/UX·게임·콘텐츠 총괄",
            "보조공학·HUSS AI 장려상 후에도 최우수팀 인터뷰·심사기준 재분석",
        ],
    },
    {
        "name": "환불원정대 팀장",
        "points": ["뉴스 분석 대시보드 방향 설정·팀 조율"],
    },
    {
        "name": "노을 밴드 동아리 운영진",
        "points": ["공연 예약·체크인 툴(bandy) 기획·구현·운영"],
    },
    {
        "name": "외부 활동",
        "points": [
            "2026.07~ LG AIMERS",
            "2026.07~09 KT 희망나눔재단 랜선나눔캠퍼스 대학생 멘토",
            "토스 현직 개발자 커피챗 → ‘문제 정의·지속 개선’ 관점 전환",
        ],
    },
]


def project_md(p: dict) -> str:
    return textwrap.dedent(
        f"""\
        # {p['name']}

        > {p['tagline']}  
        > Tier **{p['tier']}** · GitHub: {p['github']}

        ## 왜 이 문제를 선택했는가
        {p['why']}

        ## 실제 문제로 어떻게 검증했는가
        {p['validation']}

        ## 솔루션
        {p['solution']}

        ## 효과 · 정량 지표
        {p['impact']}

        ## 역할 · 역량 · 강점
        - **역할**: {p['role']}
        - **보여줄 수 있는 강점**: {p['strengths']}
        {f"- **비고**: {p['extra']}" if p.get('extra') else ''}

        ## 기술 스택
        {p['stack']}
        """
    )


def research_md(r: dict) -> str:
    return textwrap.dedent(
        f"""\
        # {r['name']} (연구형)

        > GitHub: {r['github']}

        ## 왜 이 문제를 선택했는가
        {r['why']}

        ## 실제 문제로 어떻게 검증했는가
        {r['validation']}

        ## 솔루션
        {r['solution']}

        ## 효과 · 정량 지표
        {r['impact']}

        ## 역할 · 역량 · 강점
        - **역할**: {r['role']}
        - **보여줄 수 있는 강점**: {r['strengths']}

        ## 기술 스택
        {r['stack']}
        """
    )


def write_mds() -> None:
    for p in PROJECTS:
        (PROJECTS_DIR / f"{p['slug']}.md").write_text(project_md(p), encoding="utf-8")
    for r in RESEARCH:
        (PROJECTS_DIR / f"{r['slug']}.md").write_text(research_md(r), encoding="utf-8")

    lines = [
        f"# {PROFILE['title']}",
        "",
        f"- **이름**: {PROFILE['name']}",
        f"- **소속**: {PROFILE['school']}",
        f"- **GitHub**: {PROFILE['github']}",
        f"- **성장 기록**: {PROFILE['instagram']}",
        "",
        "## 내러티브 구조",
        "1. **성적** — 꾸준한 학업 성취는 강점이다",
        "2. **학연·연구형** — 그럼 공부만 잘하는가? → 문제 검증·실험·데이터로 답한다",
        "3. **실력** — 프로젝트로 증명한다",
        "4. **인간성** — 동아리로 사람과 호흡한다",
        "5. **리더십** — 팀장·운영진·멘토로 이끈다",
        "",
        "## 1. 성적",
        f"- 누적 GPA **{GRADES['gpa']}** (환산 {GRADES['percentile']}), 취득 {GRADES['credits']}학점",
        "",
        "| 학기 | GPA | 하이라이트 |",
        "|---|---|---|",
    ]
    for sem, gpa, note in GRADES["semesters"]:
        lines.append(f"| {sem} | {gpa} | {note} |")
    lines += ["", "### 장학금"] + [f"- {s}" for s in GRADES["scholarships"]]
    lines += ["", "### 수상"] + [f"- {a}" for a in GRADES["awards"]]

    lines += [
        "",
        "## 2. 공부만 잘하는가? → 연구형·학연에 준하는 활동",
        "> 문서상 공식 ‘학연생’ 임용 기록은 확인되지 않아, **문제검증·실험·ML 연구 산출물**로 구성했습니다. 소속 랩/기간이 있으면 이 섹션을 채워 주세요.",
        "",
    ]
    for r in RESEARCH:
        lines.append(f"### [{r['name']}](projects/{r['slug']}.md)")
        lines.append(f"- {r['why']}")
        lines.append(f"- 역할: {r['role']} · 스택: {r['stack']}")
        lines.append("")

    lines += ["## 3. 실력 → 프로젝트 경험", ""]
    for tier in ("S", "A", "B"):
        lines.append(f"### Tier {tier}")
        for p in PROJECTS:
            if p["tier"] == tier:
                lines.append(f"- **[{p['name']}](projects/{p['slug']}.md)** — {p['tagline']}")
        lines.append("")

    lines += ["## 4. 인간성 → 동아리", ""]
    for c in CLUBS:
        lines.append(f"### {c['name']} ({c['role']})")
        for pt in c["points"]:
            lines.append(f"- {pt}")
        lines.append("")

    lines += ["## 5. 리더십", ""]
    for L in LEADERSHIP:
        lines.append(f"### {L['name']}")
        for pt in L["points"]:
            lines.append(f"- {pt}")
        lines.append("")

    lines += [
        "## GitHub 저장소 맵 (dlwldn4824)",
        "",
        "| Repo | 포트폴리오 매핑 |",
        "|---|---|",
        "| HOPE_organization | HOPE |",
        "| smart_icheon_care | Smart Icheon Care |",
        "| kb_AI_challenge | 답변등기 |",
        "| band | bandy |",
        "| opensource_final | Setlist Recommender |",
        "| TM-MultiLayer-MentalHealth | Mental Health 연구 |",
        "| machine_learning | 디저트 소비 ML |",
        "| noeul_boardgame | NOEUL 주루마블 |",
        "| TemiTellMe / HCI-UX / mobile_robot_temi / coshow / HCI | TEMI·HCI |",
        "| creative_py_project / open-source-fastapi-docker / Open_* | 실습 묶음 |",
        "| ik-s/iM-Ready (협업) | iM Ready |",
        "| KWwoojin/project (협업) | WJVOX |",
        "| bhw119/Cam-Kit (협업) | Cam-Kit |",
        "| PDFMathTranslate / 2026_exercise_1 / 9th_web | fork·학습용 (참고) |",
        "",
        "## 디자인 레퍼런스",
        "- KB 답변등기: ink `#26282C`, bg `#F5F6F4`, accent yellow CTA 절제",
        "- iM Shield: teal `#00BFA5` 계열",
        "- WJVOX UNIVERSAL_UI: brand-first, 한 화면 한 일, Pretendard + display font",
        "- 본 PPTX: ink + teal accent + Pretendard (위 토큰 통합)",
        "",
    ]
    (ROOT / "포트폴리오_통합.md").write_text("\n".join(lines), encoding="utf-8")


# --- PPTX helpers ------------------------------------------------------------

def _set_run_font(run, size_pt: float, bold: bool = False, color: RGBColor = INK, font: str = FONT):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    # East Asian font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def add_bg(slide, prs):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    # top accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_footer(slide, prs, page: str):
    box = slide.shapes.add_textbox(
        Inches(0.6), prs.slide_height - Inches(0.45), Inches(8), Inches(0.3)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"이지우 · dlwldn4824  ·  {page}"
    _set_run_font(run, 10, False, MUTED)
    # right accent
    box2 = slide.shapes.add_textbox(
        prs.slide_width - Inches(2.2), prs.slide_height - Inches(0.45), Inches(1.6), Inches(0.3)
    )
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = "WHY → BUILD"
    _set_run_font(run2, 10, True, ACCENT)


def add_title_block(slide, title: str, subtitle: str | None = None, y=0.35):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.2), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run_font(run, 28, True, INK)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.55), Inches(12.2), Inches(0.4))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        _set_run_font(run2, 14, False, MUTED)


def add_bullets(slide, items: list[str], left=0.6, top=1.3, width=12.2, height=5.2, size=15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        _set_run_font(run, size, False, INK)


def add_two_col(slide, left_items, right_items, top=1.4):
    add_bullets(slide, left_items, 0.6, top, 6.0, 5.0, 14)
    add_bullets(slide, right_items, 6.9, top, 6.0, 5.0, 14)


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_bg(slide, prs)
    return slide


def kpi_row(slide, kpis: list[tuple[str, str]], top=1.3):
    n = len(kpis)
    total_w = 12.2
    gap = 0.2
    w = (total_w - gap * (n - 1)) / n
    for i, (label, value) in enumerate(kpis):
        x = 0.6 + i * (w + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(1.15)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LINE
        # left accent
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(top), Inches(0.08), Inches(1.15)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT
        accent.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(top + 0.15), Inches(w - 0.3), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = value
        _set_run_font(run, 22, True, INK)
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = label
        _set_run_font(run2, 11, False, MUTED)


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    page = 0

    def stamp(slide):
        nonlocal page
        page += 1
        add_footer(slide, prs, f"{page}")

    # 1 Title
    s = new_slide(prs)
    stamp(s)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = PROFILE["title"]
    _set_run_font(run, 34, True, INK)
    sub = s.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(11.5), Inches(1.4))
    tf = sub.text_frame
    tf.word_wrap = True
    for i, line in enumerate([
        PROFILE["school"],
        f"GitHub {PROFILE['github']}  ·  {PROFILE['instagram']}",
        "성적 → 연구 → 프로젝트 → 동아리 → 리더십",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        _set_run_font(run, 14, False, ACCENT if i == 0 else MUTED)

    # 2 Agenda
    s = new_slide(prs)
    stamp(s)
    add_title_block(s, "이 포트폴리오의 질문 구조", "답을 던지기 전에, 왜를 먼저 묻습니다")
    add_bullets(
        s,
        [
            "1. 성적 — 강점 중 하나이므로 자세히 보여줍니다",
            "2. 공부만 잘하는가? / 성적 따기 쉬웠나? → 연구형·문제검증(학연에 준하는 활동)",
            "3. 실력은? → 프로젝트 경험 (문제선택→검증→솔루션→지표→역할·스택)",
            "4. 인간성은? → 동아리 (노을 운영진, 창업 동아리)",
            "5. 리더십은? → 팀장·운영진·멘토·외부 활동",
        ],
        top=1.5,
        size=18,
    )

    # 3 Grades overview
    s = new_slide(prs)
    stamp(s)
    add_title_block(s, "1. 성적 — 꾸준함이 만든 숫자", "광운대 정보융합 · 비주얼테크놀로지")
    kpi_row(
        s,
        [
            ("누적 GPA", GRADES["gpa"]),
            ("환산 점수", GRADES["percentile"]),
            ("취득 학점", GRADES["credits"]),
            ("만점 학기", "2회 (24-2, 26-1)"),
        ],
        top=1.35,
    )
    add_bullets(
        s,
        [
            "학과 특성상 상위권 경쟁이 치열하고 성적장학금 쿼터가 제한적 — 만점에도 전액이 보장되지 않음(자소서)",
            "성적은 ‘쉬워서’가 아니라 매 학기 실행 루프(수업·프로젝트·회고)의 부산물",
            "강점 축: AI/ML · 데이터 · 프로그래밍 · UX/HCI 전 구간 A+",
        ],
        top=2.8,
        size=15,
    )

    # 4 Semester table-like
    s = new_slide(prs)
    stamp(s)
    add_title_block(s, "학기별 GPA & 핵심 과목", "성적증명서 기준")
    items = [f"{sem}  ·  GPA {gpa}  —  {note}" for sem, gpa, note in GRADES["semesters"]]
    add_bullets(s, items, top=1.4, size=15)

    # 5 Scholarships & awards
    s = new_slide(prs)
    stamp(s)
    add_title_block(s, "장학금 · 수상", "학업 성취가 외부에서도 인정된 기록")
    left = ["[장학금]"] + GRADES["scholarships"]
    right = ["[수상]"] + GRADES["awards"]
    add_two_col(s, left, right, top=1.35)

    # 6 Research intro
    s = new_slide(prs)
    stamp(s)
    add_title_block(
        s,
        "2. 공부만 잘하는 거 아냐?",
        "성적 따기 쉬웠던 거 아냐? → 연구형으로 답합니다",
    )
    add_bullets(
        s,
        [
            "공식 학연생 임용 문서는 아직 미기재 — 아래는 문제검증·실험·ML 연구 산출물",
            "HOPE: 뉴스·논문·공공데이터로 Pain을 계량화한 뒤 기능을 매핑",
            "Mental Health Multi-Agent: Single LLM vs Multi-Agent+RAG 안전성 평가",
            "디저트 소비 ML: 서울시 상권 5개년 실데이터로 모델 비교",
            "공통 메시지: ‘답을 내기 전에 왜가 데이터인가’를 먼저 확인한다",
        ],
        top=1.5,
        size=16,
    )

    for r in RESEARCH:
        s = new_slide(prs)
        stamp(s)
        add_title_block(s, f"연구 · {r['name']}", r["role"])
        add_bullets(
            s,
            [
                f"왜: {r['why']}",
                f"검증: {r['validation']}",
                f"솔루션: {r['solution']}",
                f"효과: {r['impact']}",
                f"강점: {r['strengths']}",
                f"스택: {r['stack']}",
                f"GitHub: {r['github']}",
            ],
            top=1.35,
            size=14,
        )

    # Projects overview
    s = new_slide(prs)
    stamp(s)
    add_title_block(s, "3. 실력은? → 프로젝트 경험", "Tier S / A / B — 경험 누락 없이")
    s_list = [p["name"] for p in PROJECTS if p["tier"] == "S"]
    a_list = [p["name"] for p in PROJECTS if p["tier"] == "A"]
    b_list = [p["name"] for p in PROJECTS if p["tier"] == "B"]
    add_bullets(
        s,
        [
            "Tier S: " + " · ".join(s_list),
            "Tier A: " + " · ".join(a_list),
            "Tier B: " + " · ".join(b_list),
            "각 프로젝트 슬라이드 구성: 왜 → 검증 → 솔루션 → 지표 → 역할·강점 → 스택",
        ],
        top=1.5,
        size=16,
    )

    for p in PROJECTS:
        # slide 1
        s = new_slide(prs)
        stamp(s)
        add_title_block(s, f"[{p['tier']}] {p['name']}", p["tagline"])
        add_bullets(
            s,
            [
                f"왜 선택: {p['why']}",
                f"검증: {p['validation']}",
                f"솔루션: {p['solution']}",
            ],
            top=1.35,
            size=14,
        )
        # slide 2
        s = new_slide(prs)
        stamp(s)
        add_title_block(s, f"{p['name']} — 효과 · 역할 · 스택", p["github"])
        add_bullets(
            s,
            [
                f"효과/지표: {p['impact']}",
                f"역할: {p['role']}",
                f"강점: {p['strengths']}",
                f"스택: {p['stack']}",
            ]
            + ([f"비고: {p['extra']}"] if p.get("extra") else []),
            top=1.35,
            size=14,
        )

    # Clubs
    s = new_slide(prs)
    stamp(s)
    add_title_block(s, "4. 인간성은? → 동아리", "기술로 커뮤니티를 돕는 사람")
    items = []
    for c in CLUBS:
        items.append(f"{c['name']} — {c['role']}")
        items.extend(c["points"])
    add_bullets(s, items, top=1.4, size=15)

    # Leadership
    s = new_slide(prs)
    stamp(s)
    add_title_block(s, "5. 리더십", "팀장 · 운영진 · 멘토")
    items = []
    for L in LEADERSHIP:
        items.append(f"[{L['name']}]")
        items.extend(L["points"])
    add_bullets(s, items, top=1.35, size=14)

    # Closing
    s = new_slide(prs)
    stamp(s)
    add_title_block(s, "한 줄로", "프로젝트는 결과를 증명하는 과정이 아니라, 다음 프로젝트를 위한 데이터입니다")
    add_bullets(
        s,
        [
            "왜를 묻는다 → 데이터로 검증한다 → 되돌릴 수 없는 곳엔 결정론을 둔다",
            "학업(4.32) · 연구형 검증 · 제품 프로젝트 · 동아리 · 리더십을 한 루프로 돌린다",
            "연락/포트폴리오: GitHub dlwldn4824 · @due_study_archive",
            "상세 MD: 포트폴리오_통합.md + projects/*.md",
        ],
        top=2.0,
        size=16,
    )

    out = ROOT / "이지우_포트폴리오.pptx"
    prs.save(out)
    return out


def main():
    write_mds()
    out = build_pptx()
    md_count = len(list(PROJECTS_DIR.glob("*.md")))
    print(f"Wrote {md_count} project/research MDs")
    print(f"Wrote {ROOT / '포트폴리오_통합.md'}")
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
